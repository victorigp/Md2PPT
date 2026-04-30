# -*- coding: utf-8 -*-
"""
Md2PPT.py - Conversor de Markdown a PowerPoint usando una plantilla .pptx

Uso:
    python Md2PPT.py <archivo.md> <plantilla.pptx> <salida.pptx> [--debug]
"""

import argparse
import os
import re
import subprocess
import tempfile
from datetime import datetime
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


# ---------------------------------------------------------------------------
# Nombres canonicos asignados por LimpiarPlantilla.py
# ---------------------------------------------------------------------------

N_TITULO  = "MD2PPT_TITLE"
N_SUBTIT  = "MD2PPT_SUBTITLE"
N_ANTETIT = "MD2PPT_ANTETITLE"
N_CUERPO  = "MD2PPT_BODY"
N_IMAGEN  = "MD2PPT_PICTURE"
N_PIE     = "MD2PPT_FOOTER"
N_NUMERO  = "MD2PPT_SLIDE_NUMBER"
N_FECHA   = "MD2PPT_DATE"

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Regex imagen Markdown: ![alt](ruta) o ![alt](ruta "titulo")
RE_IMG = re.compile(r'^!\[.*?\]\((.+?)(?:\s+".*?")?\)\s*$')


# ---------------------------------------------------------------------------
# Helpers para clonar placeholders del layout a la slide
# ---------------------------------------------------------------------------

def _spTree(slide):
    """Devuelve el spTree de la slide o None."""
    cSld = slide._element.find(f'.//{{{NS_P}}}cSld')
    if cSld is None:
        return None
    return cSld.find(f'.//{{{NS_P}}}spTree')


def clonar_ph_desde_layout(layout, slide, idx):
    """Clona el shape con ph idx=N del layout al slide. Devuelve el clon o None."""
    ns = {"p": NS_P}
    for sp in layout._element.findall('.//p:sp', ns):
        ph_el = sp.find('.//p:ph', ns)
        if ph_el is not None and ph_el.get('idx') == str(idx):
            spTree = _spTree(slide)
            if spTree is not None:
                sp_clone = deepcopy(sp)
                spTree.append(sp_clone)
                return sp_clone
    return None


# ---------------------------------------------------------------------------
# Busqueda y escritura de placeholders
# ---------------------------------------------------------------------------

def buscar_ph(obj, nombre=None, tipo=None, idx=None):
    """Busca placeholder: nombre canonico > tipo enum > idx numerico."""
    if nombre:
        for ph in obj.placeholders:
            if ph.name == nombre:
                return ph
    if tipo is not None:
        for ph in obj.placeholders:
            try:
                if ph.placeholder_format.type == tipo:
                    return ph
            except Exception:
                continue
    if idx is not None:
        for ph in obj.placeholders:
            try:
                if ph.placeholder_format.idx == idx:
                    return ph
            except Exception:
                continue
    return None


def escribir_texto_preservando(ph, texto):
    """Escribe texto en un placeholder preservando el formato del primer run."""
    tf = ph.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = texto
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
    else:
        ph.text = texto


# ---------------------------------------------------------------------------
# Parseo de Markdown
# ---------------------------------------------------------------------------

REGEX_INLINE = re.compile(
    r'(\*\*|__)(.*?)\1'
    r'|(\*|_)(.*?)\3'
    r'|~~(.*?)~~'
    r'|`(.*?)`'
    r'|([^*_~`]+)',
    re.DOTALL
)


def parsear_inline(texto):
    """Parsea negrita, cursiva, tachado y codigo inline."""
    tokens = []
    for m in REGEX_INLINE.finditer(texto):
        if m.group(1):
            tokens.append({"texto": m.group(2), "bold": True,  "italic": False, "strike": False, "code": False})
        elif m.group(3):
            tokens.append({"texto": m.group(4), "bold": False, "italic": True,  "strike": False, "code": False})
        elif m.group(5):
            tokens.append({"texto": m.group(5), "bold": False, "italic": False, "strike": True,  "code": False})
        elif m.group(6):
            tokens.append({"texto": m.group(6), "bold": False, "italic": False, "strike": False, "code": True})
        elif m.group(7):
            tokens.append({"texto": m.group(7), "bold": False, "italic": False, "strike": False, "code": False})
    return tokens


def parsear_markdown(ruta_md):
    """Lee el MD y devuelve (lista_elementos, primer_imagen, primer_mermaid)."""
    elementos      = []
    primer_mermaid = None
    primer_imagen  = None
    en_mermaid     = False
    bloque_mermaid = []
    md_dir = os.path.dirname(os.path.abspath(ruta_md))

    with open(ruta_md, "r", encoding="utf-8") as f:
        for linea in f:
            linea = linea.rstrip("\n")

            # Bloque mermaid
            if re.match(r"^\s*```mermaid\s*$", linea):
                en_mermaid = True
                bloque_mermaid = []
                continue
            if en_mermaid and re.match(r"^\s*```\s*$", linea):
                en_mermaid = False
                if bloque_mermaid:
                    codigo = "\n".join(bloque_mermaid)
                    if primer_mermaid is None:
                        primer_mermaid = codigo
                    elementos.append({"tipo": "mermaid", "codigo": codigo})
                continue
            if en_mermaid:
                bloque_mermaid.append(linea)
                continue

            if not linea.strip():
                continue

            # Imagen Markdown: ![alt](ruta)
            m = RE_IMG.match(linea)
            if m:
                ruta_img = os.path.join(md_dir, m.group(1).strip())
                if primer_imagen is None:
                    primer_imagen = ruta_img
                elementos.append({"tipo": "imagen", "ruta": ruta_img})
                continue

            m = re.match(r"^###\s+(.+)$", linea)
            if m:
                elementos.append({"tipo": "h3", "texto": m.group(1).strip()})
                continue
            m = re.match(r"^##\s+(.+)$", linea)
            if m:
                elementos.append({"tipo": "h2", "texto": m.group(1).strip()})
                continue
            m = re.match(r"^#\s+(.+)$", linea)
            if m:
                elementos.append({"tipo": "h1", "texto": m.group(1).strip()})
                continue
            m = re.match(r"^\s*[*\-]\s+(.+)$", linea)
            if m:
                elementos.append({"tipo": "bullet", "texto": m.group(1).strip()})
                continue
            elementos.append({"tipo": "text", "texto": linea.strip()})

    return elementos, primer_imagen, primer_mermaid


def etiquetar_mermaid_secciones(elementos):
    """Anota en cada h2 la primera imagen y el primer mermaid de su seccion."""
    for i, el in enumerate(elementos):
        if el["tipo"] == "h2":
            el["primer_imagen"]  = None
            el["primer_mermaid"] = None
            for j in range(i + 1, len(elementos)):
                if elementos[j]["tipo"] in ("h1", "h2"):
                    break
                if elementos[j]["tipo"] == "imagen" and el["primer_imagen"] is None:
                    el["primer_imagen"] = elementos[j]["ruta"]
                if elementos[j]["tipo"] == "mermaid" and el["primer_mermaid"] is None:
                    el["primer_mermaid"] = elementos[j]["codigo"]
    return elementos


# ---------------------------------------------------------------------------
# Utilidades de presentacion
# ---------------------------------------------------------------------------

def obtener_mes_anio_actual():
    meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
             "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"]
    return f"{meses[datetime.now().month - 1]} {datetime.now().year}"


def renderizar_mermaid(codigo_mermaid):
    """Renderiza Mermaid a PNG. Devuelve ruta temporal o None."""
    try:
        tmp_mmd = tempfile.NamedTemporaryFile(suffix=".mmd", delete=False, mode="w", encoding="utf-8")
        tmp_mmd.write(codigo_mermaid)
        tmp_mmd.close()
        tmp_png = tmp_mmd.name.replace(".mmd", ".png")
        resultado = subprocess.run(
            ["mmdc", "-i", tmp_mmd.name, "-o", tmp_png, "-b", "transparent"],
            capture_output=True, text=True, timeout=30, shell=True
        )
        os.unlink(tmp_mmd.name)
        if resultado.returncode == 0 and os.path.exists(tmp_png):
            return tmp_png
        print(f"  [AVISO] Mermaid: {resultado.stderr}")
    except FileNotFoundError:
        print("  [AVISO] mmdc no instalado. Instalar con: npm install -g @mermaid-js/mermaid-cli")
    except Exception as e:
        print(f"  [AVISO] Error Mermaid: {e}")
    return None


def insertar_imagen_en_area(slide, ruta_png, left, top, width, height):
    """Inserta una imagen centrada (fit) dentro del area dada."""
    if not HAS_PIL or not ruta_png:
        return
    with Image.open(ruta_png) as img:
        iw, ih = img.size
    ratio = min(width / iw, height / ih)
    nw, nh = int(iw * ratio), int(ih * ratio)
    slide.shapes.add_picture(ruta_png,
                             left + (width - nw) // 2,
                             top + (height - nh) // 2,
                             nw, nh)


def anadir_parrafo_contenido(texto_frame, texto, es_bullet=False):
    """Anade un parrafo con formato inline Markdown al text_frame."""
    if texto_frame.paragraphs and not texto_frame.paragraphs[0].text:
        p = texto_frame.paragraphs[0]
    else:
        p = texto_frame.add_paragraph()

    p.text = ""
    p.level = 0

    if not es_bullet:
        pPr = p._p.get_or_add_pPr()
        for tag in [qn('a:buChar'), qn('a:buFont'), qn('a:buClr'),
                    qn('a:buSzPct'), qn('a:buSzPts'), qn('a:buSzTx'),
                    qn('a:buNone'), qn('a:buAutoNum')]:
            for el in pPr.findall(tag):
                pPr.remove(el)
        etree.SubElement(pPr, qn('a:buNone'))

    for frag in parsear_inline(texto):
        run = p.add_run()
        run.text = frag["texto"]
        if frag["bold"]:
            run.font.bold = True
        if frag["italic"]:
            run.font.italic = True
        if frag["strike"]:
            run._r.get_or_add_rPr().set('strike', 'sngStrike')
        if frag["code"]:
            run.font.bold = False
            run.font.italic = False
            rPr = run._r.get_or_add_rPr()
            latin = etree.SubElement(rPr, qn('a:latin'))
            latin.set('typeface', 'Courier New')


# ---------------------------------------------------------------------------
# Creacion de slides de contenido
# ---------------------------------------------------------------------------

def construir_mapa_idx_layout(layout):
    """
    Lee los nombres MD2PPT_* del layout anotado por LimpiarPlantilla.py
    y devuelve dict nombre -> idx.
    Si el layout no fue anotado devuelve dict vacio (se usaran fallbacks).
    """
    mapa = {}
    for ph in layout.placeholders:
        if ph.name and ph.name.startswith("MD2PPT_"):
            mapa[ph.name] = ph.placeholder_format.idx
    return mapa


def crear_slide_contenido(prs, layout, titulo_seccion, titulo_presentacion,
                          subtitulo=None, mapa_idx=None, company=""):
    """
    Crea y configura una slide de contenido.
    mapa_idx: dict nombre_canonico -> idx del layout anotado.
    Devuelve (slide, texto_frame_del_cuerpo).
    """
    slide = prs.slides.add_slide(layout)

    # Mapa idx -> placeholder de la slide (una sola pasada, fiable)
    def ph_map():
        m = {}
        for ph in slide.placeholders:
            try:
                m[ph.placeholder_format.idx] = ph
            except Exception:
                continue
        return m

    phd = ph_map()

    # Resolver idx de cada rol: mapa anotado > fallback estandar
    mi = mapa_idx or {}
    idx_titulo  = mi.get(N_TITULO,  0)
    idx_subtit  = mi.get(N_SUBTIT,  14)
    idx_antetit = mi.get(N_ANTETIT, 16)
    idx_cuerpo  = mi.get(N_CUERPO,  15)
    idx_pie     = mi.get(N_PIE,     17)
    idx_numero  = mi.get(N_NUMERO,  18)

    # Pie y numero de pagina: clonar desde layout si no los instancio pptx
    for idx_extra in (idx_pie, idx_numero):
        if idx_extra not in phd:
            clonar_ph_desde_layout(layout, slide, idx_extra)
            phd = ph_map()

    # Titulo
    if idx_titulo in phd:
        escribir_texto_preservando(phd[idx_titulo], titulo_seccion)

    # Antetitulo: eliminar siempre
    if idx_antetit in phd:
        try:
            phd[idx_antetit]._element.getparent().remove(phd[idx_antetit]._element)
        except Exception:
            pass

    # Subtitulo: escribir o eliminar
    if idx_subtit in phd:
        if subtitulo:
            escribir_texto_preservando(phd[idx_subtit], subtitulo)
        else:
            try:
                phd[idx_subtit]._element.getparent().remove(phd[idx_subtit]._element)
            except Exception:
                pass

    # Pie
    if idx_pie in phd:
        texto_pie = f"{company}  /  {titulo_presentacion}" if company else titulo_presentacion
        escribir_texto_preservando(phd[idx_pie], texto_pie)

    # Cuerpo
    tf_contenido = None
    if idx_cuerpo in phd:
        tf_contenido = phd[idx_cuerpo].text_frame
        tf_contenido.clear()

    return slide, tf_contenido


# ---------------------------------------------------------------------------
# Helpers para insertar imagen en slide de contenido
# ---------------------------------------------------------------------------

def _area_cuerpo(slide_actual, mapa_idx_contenido):
    """Devuelve (left, top, width, height) del placeholder de cuerpo o fallback."""
    idx_cuerpo = mapa_idx_contenido.get(N_CUERPO, 15)
    for ph in slide_actual.placeholders:
        try:
            if ph.placeholder_format.idx == idx_cuerpo:
                return ph.left, ph.top, ph.width, ph.height
        except Exception:
            continue
    return Inches(0.5), Inches(2), Inches(9), Inches(4)


# ---------------------------------------------------------------------------
# Generacion principal
# ---------------------------------------------------------------------------

def generar_presentacion(elementos, ruta_plantilla, ruta_salida,
                         primer_imagen=None, mermaid_code=None, company=""):
    """Genera la presentacion PowerPoint."""
    elementos = etiquetar_mermaid_secciones(elementos)
    prs = Presentation(ruta_plantilla)

    # Limpiar slides existentes
    xml_slides = prs.slides._sldIdLst
    while len(xml_slides) > 0:
        prs.part.drop_rel(xml_slides[0].rId)
        del xml_slides[0]

    if len(prs.slide_layouts) < 4:
        raise ValueError(
            f"La plantilla tiene {len(prs.slide_layouts)} layouts. "
            f"Ejecuta LimpiarPlantilla.py primero."
        )

    layout_portada   = prs.slide_layouts[0]
    layout_seccion   = prs.slide_layouts[1]
    layout_contenido = prs.slide_layouts[2]
    layout_cierre    = prs.slide_layouts[3]

    # Mapa nombre->idx del layout de contenido, leido una sola vez
    mapa_idx_contenido = construir_mapa_idx_layout(layout_contenido)

    # Imagen para la portada: preferir imagen directa, fallback a mermaid renderizado
    ruta_portada_png   = None
    mermaid_portada_tmp = None  # solo para limpieza si fue generado
    if primer_imagen and os.path.exists(primer_imagen):
        ruta_portada_png = primer_imagen
    elif mermaid_code:
        mermaid_portada_tmp = renderizar_mermaid(mermaid_code)
        ruta_portada_png    = mermaid_portada_tmp

    titulo_presentacion   = ""
    contador_seccion      = 0
    titulo_seccion_actual = ""
    slide_seccion_actual  = None
    subtitulo_seccion_ph  = None
    texto_frame_actual    = None
    seccion_pendiente     = False

    def nueva_slide_contenido(subtitulo=None):
        nonlocal texto_frame_actual, seccion_pendiente
        _, tf = crear_slide_contenido(
            prs, layout_contenido,
            titulo_seccion_actual, titulo_presentacion,
            subtitulo=subtitulo,
            mapa_idx=mapa_idx_contenido,
            company=company
        )
        texto_frame_actual = tf
        seccion_pendiente  = False

    def asegurar_slide_contenido():
        if seccion_pendiente:
            nueva_slide_contenido()

    for elemento in elementos:

        # ---------------------------------------------------------------- H1
        if elemento["tipo"] == "h1":
            titulo_presentacion = elemento["texto"]
            slide = prs.slides.add_slide(layout_portada)

            ph_t = buscar_ph(slide, N_TITULO, PP_PLACEHOLDER.TITLE, 0) \
                   or buscar_ph(slide, tipo=PP_PLACEHOLDER.CENTER_TITLE)
            if ph_t:
                ph_t.text = titulo_presentacion

            ph_s = buscar_ph(slide, N_SUBTIT, PP_PLACEHOLDER.SUBTITLE, 1)
            if ph_s:
                ph_s.text = ""

            ph_f = buscar_ph(slide, N_FECHA, PP_PLACEHOLDER.DATE)
            if ph_f:
                ph_f.text = obtener_mes_anio_actual()
            else:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx > 1 and \
                           ph.placeholder_format.type != PP_PLACEHOLDER.PICTURE:
                            ph.text = obtener_mes_anio_actual()
                            break
                    except Exception:
                        continue

            if ruta_portada_png and HAS_PIL:
                ph_img = buscar_ph(slide, N_IMAGEN, PP_PLACEHOLDER.PICTURE)
                if ph_img:
                    insertar_imagen_en_area(slide, ruta_portada_png,
                                           ph_img.left, ph_img.top, ph_img.width, ph_img.height)
                    ph_img._element.getparent().remove(ph_img._element)

        # ---------------------------------------------------------------- H3
        elif elemento["tipo"] == "h3":
            if slide_seccion_actual is not None and subtitulo_seccion_ph is not None:
                t = subtitulo_seccion_ph.text
                subtitulo_seccion_ph.text = (t + "\n" + elemento["texto"]) if t else elemento["texto"]
            nueva_slide_contenido(subtitulo=elemento["texto"])

        # ---------------------------------------------------------------- H2
        elif elemento["tipo"] == "h2":
            contador_seccion += 1
            titulo_seccion_actual = elemento["texto"]

            slide_seccion_actual = prs.slides.add_slide(layout_seccion)

            ph_t = buscar_ph(slide_seccion_actual, N_TITULO, PP_PLACEHOLDER.TITLE, 0)
            if ph_t:
                ph_t.text = elemento["texto"].upper()

            subtitulo_seccion_ph = buscar_ph(slide_seccion_actual, N_SUBTIT, PP_PLACEHOLDER.SUBTITLE, 1)
            if subtitulo_seccion_ph:
                subtitulo_seccion_ph.text = ""

            ph_num = buscar_ph(slide_seccion_actual, idx=10)
            if ph_num:
                escribir_texto_preservando(ph_num, f"{contador_seccion:02d}")

            # Imagen de seccion: preferir imagen directa, fallback a mermaid
            ruta_sec       = None
            mermaid_sec_tmp = None
            img_sec = elemento.get("primer_imagen")
            if img_sec and os.path.exists(img_sec):
                ruta_sec = img_sec
            else:
                mermaid_sec = elemento.get("primer_mermaid")
                if mermaid_sec:
                    mermaid_sec_tmp = renderizar_mermaid(mermaid_sec)
                    ruta_sec = mermaid_sec_tmp

            if ruta_sec:
                try:
                    ph_img = buscar_ph(slide_seccion_actual, N_IMAGEN, PP_PLACEHOLDER.PICTURE)
                    if ph_img:
                        ar_l, ar_t, ar_w, ar_h = ph_img.left, ph_img.top, ph_img.width, ph_img.height
                    else:
                        ar_l, ar_t, ar_w, ar_h = Inches(6), Inches(1), Inches(3.5), Inches(5)
                    insertar_imagen_en_area(slide_seccion_actual, ruta_sec, ar_l, ar_t, ar_w, ar_h)
                except Exception as e:
                    print(f"  [AVISO] Imagen en portada seccion: {e}")
                finally:
                    if mermaid_sec_tmp and os.path.exists(mermaid_sec_tmp):
                        os.unlink(mermaid_sec_tmp)

            seccion_pendiente  = True
            texto_frame_actual = None

        # ---------------------------------------------------------------- BULLET
        elif elemento["tipo"] == "bullet":
            if texto_frame_actual is not None or seccion_pendiente:
                asegurar_slide_contenido()
                anadir_parrafo_contenido(texto_frame_actual, elemento["texto"], es_bullet=True)

        # ---------------------------------------------------------------- TEXT
        elif elemento["tipo"] == "text":
            if texto_frame_actual is not None or seccion_pendiente:
                asegurar_slide_contenido()
                anadir_parrafo_contenido(texto_frame_actual, elemento["texto"], es_bullet=False)

        # ---------------------------------------------------------------- IMAGEN
        elif elemento["tipo"] == "imagen":
            if texto_frame_actual is not None or seccion_pendiente:
                asegurar_slide_contenido()
                ruta_img = elemento["ruta"]
                if os.path.exists(ruta_img):
                    try:
                        slide_actual = prs.slides[-1]
                        ar_l, ar_t, ar_w, ar_h = _area_cuerpo(slide_actual, mapa_idx_contenido)
                        insertar_imagen_en_area(slide_actual, ruta_img, ar_l, ar_t, ar_w, ar_h)
                    except Exception as e:
                        print(f"  [AVISO] Imagen en contenido: {e}")
                else:
                    print(f"  [AVISO] Imagen no encontrada: {ruta_img}")

        # ---------------------------------------------------------------- MERMAID
        elif elemento["tipo"] == "mermaid":
            if texto_frame_actual is not None or seccion_pendiente:
                asegurar_slide_contenido()
                ruta_png = renderizar_mermaid(elemento["codigo"])
                if ruta_png:
                    try:
                        slide_actual = prs.slides[-1]
                        ar_l, ar_t, ar_w, ar_h = _area_cuerpo(slide_actual, mapa_idx_contenido)
                        insertar_imagen_en_area(slide_actual, ruta_png, ar_l, ar_t, ar_w, ar_h)
                    except Exception as e:
                        print(f"  [AVISO] Mermaid en contenido: {e}")
                    finally:
                        if os.path.exists(ruta_png):
                            os.unlink(ruta_png)

    prs.slides.add_slide(layout_cierre)
    prs.save(ruta_salida)
    print(f"  Presentacion generada correctamente: {ruta_salida}")

    if mermaid_portada_tmp and os.path.exists(mermaid_portada_tmp):
        os.unlink(mermaid_portada_tmp)


# ---------------------------------------------------------------------------
# Entrada
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Convierte Markdown a PowerPoint.")
    parser.add_argument("entrada",   help="Archivo Markdown (.md)")
    parser.add_argument("plantilla", help="Plantilla PowerPoint (.pptx)")
    parser.add_argument("salida",    help="Archivo de salida (.pptx)")
    parser.add_argument("--company", default="", help="Nombre de empresa para el pie")
    parser.add_argument("--debug",   action="store_true",
                        help="Muestra layouts y placeholders de la plantilla")
    args = parser.parse_args()

    if args.debug:
        prs = Presentation(args.plantilla)
        print(f"\n  === LAYOUTS ({args.plantilla}) ===\n")
        for i, layout in enumerate(prs.slide_layouts):
            print(f"  Layout {i}: '{layout.name}'")
            for ph in layout.placeholders:
                print(f"    idx={ph.placeholder_format.idx} tipo={ph.placeholder_format.type} "
                      f"nombre='{ph.name}'")
            print()
        print("  === SLIDES EXISTENTES ===\n")
        for i, slide in enumerate(prs.slides):
            print(f"  Slide {i}: layout='{slide.slide_layout.name}'")
            for ph in slide.placeholders:
                print(f"    idx={ph.placeholder_format.idx} tipo={ph.placeholder_format.type} "
                      f"texto='{(ph.text or '')[:40]}'")
            print()
        return

    elementos, primer_imagen, primer_mermaid = parsear_markdown(args.entrada)
    if not elementos:
        print("  El archivo Markdown esta vacio o no contiene elementos reconocidos.")
        return

    try:
        generar_presentacion(elementos, args.plantilla, args.salida,
                             primer_imagen=primer_imagen,
                             mermaid_code=primer_mermaid,
                             company=args.company)
    except Exception as e:
        print(f"  Error: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
