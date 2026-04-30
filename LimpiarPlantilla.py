# -*- coding: utf-8 -*-
"""
LimpiarPlantilla.py - Limpia una plantilla .pptx eliminando layouts huerfanos
y anota los placeholders con nombres canonicos para que Md2PPT.py los detecte
dinamicamente sin depender de indices fijos.

Flujo de trabajo:
    1. Abrir la plantilla en PowerPoint y dejar solo las 4 diapositivas necesarias:
       Portada, Separata/Seccion, Contenido, Cierre.
    2. Ejecutar este script.
    3. Usar la plantilla limpia resultante con Md2PPT.py.

Uso:
    python LimpiarPlantilla.py <plantilla.pptx> <plantilla_limpia.pptx>
"""

import argparse
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


# ---------------------------------------------------------------------------
# Nombres canonicos que Md2PPT.py usara para buscar placeholders
# ---------------------------------------------------------------------------

NOMBRE_TITULO     = "MD2PPT_TITLE"
NOMBRE_SUBTITULO  = "MD2PPT_SUBTITLE"
NOMBRE_ANTETITULO = "MD2PPT_ANTETITLE"
NOMBRE_CUERPO     = "MD2PPT_BODY"
NOMBRE_IMAGEN     = "MD2PPT_PICTURE"
NOMBRE_PIE        = "MD2PPT_FOOTER"
NOMBRE_NUMERO     = "MD2PPT_SLIDE_NUMBER"
NOMBRE_FECHA      = "MD2PPT_DATE"


# ---------------------------------------------------------------------------
# Anotacion de placeholders
# ---------------------------------------------------------------------------

def anotar_placeholders_layout(layout, idx_layout):
    """
    Asigna nombres canonicos MD2PPT_* a los placeholders del layout.

    Para placeholders tipo=2 (BODY), que en plantillas corporativas pueden ser
    cuerpo, subtitulo y antetitulo a la vez, los distingue por area y posicion:
      - El de mayor area                          -> MD2PPT_BODY
      - Si queda 1 adicional                      -> MD2PPT_SUBTITLE
      - Si quedan 2+: el mas alto (menor top)     -> MD2PPT_ANTETITLE
                      el siguiente                -> MD2PPT_SUBTITLE
    """
    anotados = []

    # Recopilar todos los placeholders con su tipo, area y posicion
    todos = []
    for ph in layout.placeholders:
        try:
            tipo = ph.placeholder_format.type
        except Exception:
            continue
        try:
            area = (ph.width or 0) * (ph.height or 0)
            top  = ph.top if ph.top is not None else 0
        except Exception:
            area, top = 0, 0
        todos.append((ph, tipo, area, top))

    # Separar tipo BODY (tipo=2) del resto
    tipo2 = [(ph, area, top) for ph, tipo, area, top in todos
             if tipo == PP_PLACEHOLDER.BODY]
    otros = [(ph, tipo) for ph, tipo, area, top in todos
             if tipo != PP_PLACEHOLDER.BODY]

    # Asignar nombres a tipo=2 segun area y posicion
    if tipo2:
        por_area = sorted(tipo2, key=lambda x: x[1], reverse=True)
        # El de mayor area es el cuerpo
        ph_cuerpo = por_area[0][0]
        ph_cuerpo.name = NOMBRE_CUERPO
        anotados.append(f"idx={ph_cuerpo.placeholder_format.idx} -> {NOMBRE_CUERPO}")

        resto = por_area[1:]
        if len(resto) == 1:
            # Solo uno mas: es el subtitulo
            resto[0][0].name = NOMBRE_SUBTITULO
            anotados.append(f"idx={resto[0][0].placeholder_format.idx} -> {NOMBRE_SUBTITULO}")
        elif len(resto) >= 2:
            # Ordenar por top: el mas alto (menor top) es el antetitulo
            por_top = sorted(resto, key=lambda x: x[2])
            por_top[0][0].name = NOMBRE_ANTETITULO
            anotados.append(f"idx={por_top[0][0].placeholder_format.idx} -> {NOMBRE_ANTETITULO}")
            por_top[1][0].name = NOMBRE_SUBTITULO
            anotados.append(f"idx={por_top[1][0].placeholder_format.idx} -> {NOMBRE_SUBTITULO}")
            # Los demas (si hubiera) se dejan sin nombre canonico

    # Asignar nombres al resto segun tipo enum y nombre original
    for ph, tipo in otros:
        nombre = None
        if tipo in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            nombre = NOMBRE_TITULO
        elif tipo == PP_PLACEHOLDER.SUBTITLE:
            nombre = NOMBRE_SUBTITULO
        elif tipo == PP_PLACEHOLDER.PICTURE:
            nombre = NOMBRE_IMAGEN
        elif tipo == PP_PLACEHOLDER.FOOTER:
            nombre = NOMBRE_PIE
        elif tipo == PP_PLACEHOLDER.SLIDE_NUMBER:
            nombre = NOMBRE_NUMERO
        elif tipo == PP_PLACEHOLDER.DATE:
            nombre = NOMBRE_FECHA
        else:
            nl = (ph.name or "").lower()
            if "footer" in nl or "pie" in nl:
                nombre = NOMBRE_PIE
            elif "subtitle" in nl or "subt" in nl:
                nombre = NOMBRE_SUBTITULO
            elif "title" in nl or "tulo" in nl:
                nombre = NOMBRE_TITULO
            elif "picture" in nl or "image" in nl:
                nombre = NOMBRE_IMAGEN
            elif "number" in nl or "num" in nl:
                nombre = NOMBRE_NUMERO
            elif "date" in nl or "fecha" in nl:
                nombre = NOMBRE_FECHA
        if nombre:
            ph.name = nombre
            anotados.append(f"idx={ph.placeholder_format.idx} -> {nombre}")

    if anotados:
        print(f"    Layout {idx_layout} ('{layout.name}'):")
        for a in anotados:
            print(f"      {a}")


def anotar_todos_los_layouts(prs):
    """Anota todos los layouts de la presentacion con nombres canonicos."""
    for i, layout in enumerate(prs.slide_layouts):
        anotar_placeholders_layout(layout, i)


# ---------------------------------------------------------------------------
# Limpieza de layouts huerfanos
# ---------------------------------------------------------------------------

def obtener_layouts_en_uso(prs):
    """Devuelve el conjunto de id() de los layouts usados por las slides visibles."""
    return {id(slide.slide_layout._element) for slide in prs.slides}


def limpiar_plantilla(ruta_entrada, ruta_salida):
    prs = Presentation(ruta_entrada)

    total_layouts = len(prs.slide_layouts)
    total_slides  = len(prs.slides)

    print(f"\n  Plantilla: {ruta_entrada}")
    print(f"  Diapositivas visibles: {total_slides}")
    print(f"  Layouts totales:       {total_layouts}")

    if total_slides == 0:
        print("\n  [ERROR] La plantilla no tiene diapositivas visibles.")
        print("          Deja las diapositivas necesarias antes de ejecutar este script.")
        sys.exit(1)

    elementos_en_uso = obtener_layouts_en_uso(prs)

    print(f"\n  Layouts en uso ({len(elementos_en_uso)}):")
    for i, slide in enumerate(prs.slides):
        print(f"    Slide {i}: '{slide.slide_layout.name}'")

    # Eliminar layouts huerfanos
    slide_master = prs.slide_master
    sldLayoutIdLst = slide_master._element.find(
        './/{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutIdLst'
    )
    if sldLayoutIdLst is None:
        print("\n  [ERROR] No se encontro la lista de layouts en el master.")
        sys.exit(1)

    rels = slide_master.part.rels
    rids_a_eliminar = [
        rId for rId, rel in rels.items()
        if 'slideLayout' in rel.reltype
        and id(rel.target_part._element) not in elementos_en_uso
    ]

    eliminados = 0
    for el in list(sldLayoutIdLst):
        rId = el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId in rids_a_eliminar:
            sldLayoutIdLst.remove(el)
            slide_master.part.drop_rel(rId)
            eliminados += 1

    print(f"\n  Layouts eliminados:  {eliminados}")
    print(f"  Layouts conservados: {len(elementos_en_uso)}")

    print(f"\n  Orden final de layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"    Layout {i}: '{layout.name}'")

    print(f"\n  Anotando placeholders...")
    anotar_todos_los_layouts(prs)

    prs.save(ruta_salida)
    print(f"\n  Plantilla limpia guardada en: {ruta_salida}")


# ---------------------------------------------------------------------------
# Entrada
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Limpia layouts huerfanos de una plantilla .pptx y anota sus placeholders."
    )
    parser.add_argument("entrada", help="Plantilla con las diapositivas deseadas (.pptx)")
    parser.add_argument("salida",  help="Plantilla limpia resultante (.pptx)")
    args = parser.parse_args()

    if not os.path.exists(args.entrada):
        print(f"\n  [ERROR] No se encontro: {args.entrada}")
        sys.exit(1)
    if os.path.abspath(args.entrada) == os.path.abspath(args.salida):
        print("\n  [ERROR] Entrada y salida no pueden ser el mismo archivo.")
        sys.exit(1)

    limpiar_plantilla(args.entrada, args.salida)


if __name__ == "__main__":
    main()
